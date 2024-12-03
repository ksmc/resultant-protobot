import re
import io
import os
import fitz
import docx2txt
import pandas as pd
from PIL import Image
import streamlit as st
# from docx import Document
from pptx import Presentation
from llama_index.core import (
	Settings,
	VectorStoreIndex,
	ServiceContext,
	Document,
	SimpleDirectoryReader,
	StorageContext,
	load_index_from_storage
)

# Function to read contents of a pdf file
def read_pdf(bytes_data):
	text = ""
	if isinstance(bytes_data, bytes):
		doc = fitz.open(stream = bytes_data, filetype = "pdf")
	else:
		doc = fitz.open(bytes_data, filetype = "pdf")
	for page in doc:
		text += page.get_text()
	return text

# Function to read contents of a docx file
def read_docx(bytes_data):
	text = docx2txt.process(io.BytesIO(bytes_data))
	return text

# Function to read contents of a ppt file
def read_ppt(bytes_data):
	ppt = Presentation(io.BytesIO(bytes_data))
	text = ""
	for slide in ppt.slides:
		for shape in slide.shapes:
			if hasattr(shape, "text"):
				text += shape.text
	return text

# Function to read contents of a txt file
def read_txt(bytes_data):
	if isinstance(bytes_data, bytes):
		text = bytes_data
	else:
		with open(bytes_data, 'r') as file:
			text = file.read()
	return text

# Function to read contents of a csv file
def read_csv(bytes_data): 
	text = pd.read_csv(io.BytesIO(bytes_data), on_bad_lines = 'skip').to_string()
	return text

# Function to read contents of a excel file
def read_excel(bytes_data):
	text = pd.read_excel(io.BytesIO(bytes_data)).to_string()
	return text

# Function to read contents of a image file
# def process_image(uploaded_file):
# 	retunr None

@st.cache_resource(show_spinner = False)	
def load_data():
	# with st.spinner(text="Loading and indexing docs – hang tight!"):
	PERSIST_DIR = "./index/02-21-small"
	# PERSIST_DIR = "./index/02-22-large"
	if not os.path.exists(PERSIST_DIR):
		reader = SimpleDirectoryReader(input_dir= PERSIST_DIR, recursive=True, required_exts=[".pdf", ".doc", ".docx", "txt", "pptx"])
		docs = reader.load_data(num_workers = 2)
		index = VectorStoreIndex.from_documents(docs, llm = Settings.llm, embed_model = Settings.embed_model, show_progress = True, num_workers = 2)
		# store index
		index.storage_context.persist(persist_dir = PERSIST_DIR)
	else:
		# load the existing index
		storage_context = StorageContext.from_defaults(persist_dir=PERSIST_DIR)
		index = load_index_from_storage(storage_context)
	return index